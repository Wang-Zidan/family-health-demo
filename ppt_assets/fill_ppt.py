# -*- coding: utf-8 -*-
"""把 Agent Infra 初赛方案 PPT 模板填上「家庭健康 AI 管家」内容。"""
import os
from pptx import Presentation
from pptx.util import Inches

TEMPLATE = "D:/学习资料/Agent Infra初赛方案PPT框架模板.pptx"
ASSETS = os.path.dirname(os.path.abspath(__file__))
OUTPUT = "C:/Users/王子丹真人/WorkBuddy/2026-08-03-14-50-19/家庭健康AI管家_AgentInfra初赛方案.pptx"

def set_text(shape, text):
    """替换 shape 内所有文本，保留多段能力。"""
    tf = shape.text_frame
    tf.clear()
    if isinstance(text, str):
        lines = text.split('\n')
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line

def find_by_text(slide, marker):
    for shape in slide.shapes:
        if shape.has_text_frame and marker in shape.text_frame.text:
            return shape
    return None

def replace_pictures(slide, image_paths):
    """按 (top, left) 排序所有 picture 形状，依次替换为 image_paths。"""
    pics = [(sh.top, sh.left, sh.width, sh.height, sh) for sh in slide.shapes if sh.shape_type == 13]
    pics.sort(key=lambda x: (x[0], x[1]))
    if len(pics) != len(image_paths):
        print(f"  warning: {len(pics)} pictures but {len(image_paths)} images on slide")
    for i, (top, left, w, h, sh) in enumerate(pics):
        if i >= len(image_paths):
            break
        path = image_paths[i]
        if path is None:
            continue
        path = os.path.join(ASSETS, path)
        sp = sh._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(path, left, top, width=w, height=h)

prs = Presentation(TEMPLATE)

# ----- Slide 1 封面 -----
sh = find_by_text(prs.slides[0], "Agent Infra 新智基座")
if sh: set_text(sh, "家庭健康 AI 管家 · Agent Infra 初赛方案")

# ----- Slide 2 P0 一页纸速览 -----
slide2 = prs.slides[1]
repl2 = {
    "【在此填写项目名称】": "家庭健康 AI 管家",
    "【描述真实场景与核心痛点】": "中老年/慢病家庭：看不懂报告、记不住用药、不会就医准备",
    "【概述端到端解决方案】": "5 个 Agent 协作：翻译报告、管理用药、营养、就医、趋势",
    "【列 1–2 个关键差异化优势】": "持续健康闭环 + 7 条医疗安全红线 + 高风险人工确认",
    "【说明复用与迁移价值】": "可迁移慢病/母婴/养老场景；复用 4 个自研 Skill",
    "【说明当前完成度与里程碑】": "work team 已就绪，10 个工具实测通过",
}
for marker, txt in repl2.items():
    sh = find_by_text(slide2, marker)
    if sh: set_text(sh, txt)

# ----- Slide 5 场景与价值 -----
s5 = prs.slides[4]
sh = find_by_text(s5, "建议覆盖目标用户与核心痛点")
if sh: set_text(sh, "面向中老年慢病家庭，把「单次报告翻译」升级为「持续健康守护闭环」，实现报告可懂、用药安全、就医有备、趋势可知。")
sh = find_by_text(s5, "（示例）")
if sh: set_text(sh, "真实痛点场景")
replace_pictures(s5, ["scene_pain.png", "scene_value.png"])

# ----- Slide 7 方案总览 -----
s7 = prs.slides[6]
sh = find_by_text(s7, "建议用一张架构图")
if sh: set_text(sh, "以进度管家为编排中枢，5 个 Worker 按并行+串行依赖协作，调用 Mock Tool Server 与 MCP/RAG 能力底座，输出多模态健康管理服务。")
replace_pictures(s7, ["arch.png"])

# ----- Slide 9 多 Agent 协同设计 -----
s9 = prs.slides[8]
sh = find_by_text(s9, "建议覆盖 Agent 分工")
if sh: set_text(sh, "5 个 Agent 分三层：并行层（病历翻译官 + 健康数据分析师）→ 串行依赖链（用药管理师→营养顾问→就医导航员）→ 进度管家全局编排。")
replace_pictures(s9, ["agents_flow.png", "agents_state.png"])

# ----- Slide 11 Skill 工程体系 -----
s11 = prs.slides[10]
sh = find_by_text(s11, "建议覆盖 Skill 清单")
if sh: set_text(sh, "本赛题必选项：4 个核心 Skill 对应 5 个 Worker，依赖 10 个 Mock Tool；输入/输出/依赖/失败处理全部显式定义，可直接迁移复用。")
replace_pictures(s11, ["skills_map.png"])

# ----- Slide 13 工程落地、运行验证与安全可审计 -----
s13 = prs.slides[12]
sh = find_by_text(s13, "建议覆盖可运行性")
if sh: set_text(sh, "work team 已本地实测通过 10 个工具调用；全链路记录 Agent 决策与置信度，7 条安全红线 + 人类审批 + 数据脱敏 + 急诊 120 协议。")
replace_pictures(s13, ["sec_run.png", "sec_obs.png", "sec_gov.png", "sec_cloud.png"])

# ----- Slide 15 开放/开源计划 -----
s15 = prs.slides[14]
sh = find_by_text(s15, "建议覆盖可复用成果")
if sh: set_text(sh, "开源 work team YAML + mock 工具服务器 + 接口契约文档；协议 MIT；复用官方/自研 Skill 共 4 个。")

# ----- Slide 17 落地计划与进展 -----
s17 = prs.slides[16]
sh = find_by_text(s17, "建议覆盖当前进展")
if sh: set_text(sh, "当前：work team 就绪、工具实测通过；计划 6 周内完成 MVP、Web 界面、MCP 接入与演示视频。")
replace_pictures(s17, ["progress_timeline.png", "risk_ctrl.png"])

# ----- Slide 19 团队介绍 -----
s19 = prs.slides[18]
sh = find_by_text(s19, "可从以下方面介绍团队基本情况")
if sh:
    set_text(sh, [
        "成员背景：",
        "· 王子丹，大数据专业学生，擅长「复杂信息→人话」的翻译与产品设计。",
        "",
        "核心技能：",
        "· 多 Agent 协同设计、Skill 工程化、Python 原型开发与验证。",
        "",
        "团队分工：",
        "· 产品/方案设计 + 多 Agent 编排 + Skill 与工具开发 + 演示落地。",
        "",
        "过往成果：",
        "· 已自主开发 8 个 WorkBuddy Skill，覆盖磁盘整理、术语翻译、电商评价、",
        "  学习路径融合、Markdown 抽取、医疗素养、焦点锚定、工作日志聚合。",
        "",
        "作品集：个人 ~/.workbuddy/skills/ 目录下 8 个 Skill 已本地运行。",
    ])

prs.save(OUTPUT)
print("saved", OUTPUT)
